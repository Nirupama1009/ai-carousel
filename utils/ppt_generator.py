from pptx import Presentation

def save_slides_as_pptx(slides, filename="output.pptx"):
    prs = Presentation()
    for slide_text in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Slide"
        content.text = slide_text
    prs.save(filename)
    return filename
